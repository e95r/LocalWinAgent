"""Интеграционные проверки чтения и записи файловых форматов."""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from tools.files import FileManager


@pytest.fixture()
def file_manager(tmp_path: Path) -> FileManager:
    return FileManager([str(tmp_path)])


def test_create_text_file_with_content(file_manager: FileManager, tmp_path: Path) -> None:
    target = tmp_path / "notes"
    result = file_manager.create_file(str(target), content="Привет мир", kind="текст", confirmed=True)
    created = Path(result["path"])
    assert result["ok"] is True
    assert created.suffix == ".txt"
    assert created.read_text(encoding="utf-8") == "Привет мир"
    assert result["size"] > 0


def test_create_and_edit_word_document(file_manager: FileManager, tmp_path: Path) -> None:
    doc_target = tmp_path / "report"
    create_result = file_manager.create_file(str(doc_target), content="Первый абзац", kind="word", confirmed=True)
    path = Path(create_result["path"])
    document = Document(str(path))
    assert any(paragraph.text == "Первый абзац" for paragraph in document.paragraphs)

    edit_result = file_manager.edit_word(str(path), "Второй абзац", confirmed=True)
    assert edit_result["ok"] is True
    document = Document(str(path))
    assert document.paragraphs[-1].text == "Второй абзац"


def test_create_excel_and_edit_cell(file_manager: FileManager, tmp_path: Path) -> None:
    excel_target = tmp_path / "table"
    create_result = file_manager.create_file(str(excel_target), kind="excel", confirmed=True)
    path = Path(create_result["path"])
    edit_result = file_manager.edit_excel(str(path), "A1", "hello", confirmed=True)
    assert edit_result["ok"] is True
    workbook = load_workbook(str(path))
    sheet = workbook.active
    assert sheet["A1"].value == "hello"


def test_create_pptx_and_add_slide(file_manager: FileManager, tmp_path: Path) -> None:
    pptx_target = tmp_path / "slides"
    create_result = file_manager.create_file(str(pptx_target), kind="pptx", confirmed=True)
    path = Path(create_result["path"])
    edit_result = file_manager.edit_pptx(str(path), "Название", confirmed=True)
    assert edit_result["ok"] is True
    presentation = Presentation(str(path))
    assert len(presentation.slides) >= 1


def test_requires_confirmation_outside_allowlist(file_manager: FileManager, tmp_path: Path) -> None:
    outside = tmp_path.parent / "outside.txt"
    result = file_manager.write_text(str(outside), "data", confirmed=False)
    assert result["ok"] is False
    assert result.get("requires_confirmation") is True
    assert "подтверждение" in result["error"].lower()
    assert not outside.exists()
