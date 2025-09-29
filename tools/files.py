"""Инструменты работы с файлами для LocalWinAgent."""

from __future__ import annotations

import logging
import os
import platform
import shutil
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Iterable, Optional
from xml.sax.saxutils import escape
from zipfile import ZIP_DEFLATED, ZipFile

import config
from tools.apps import open_with_shell
from tools.exceptions import (
    EverythingNotInstalledError,
    FileOperationError,
    NotAllowedPathError,
)

try:  # pragma: no cover - импорт проверяется в рантайме
    from docx import Document  # type: ignore
    from docx.opc.exceptions import PackageNotFoundError  # type: ignore
except Exception:  # pragma: no cover - отсутствие зависимостей
    Document = None  # type: ignore[assignment]

    class PackageNotFoundError(Exception):
        """Заглушка для обработки повреждённых документов."""

try:  # pragma: no cover
    from openpyxl import Workbook, load_workbook
except ImportError as exc:  # pragma: no cover
    raise EverythingNotInstalledError("Требуется установить пакет openpyxl") from exc

try:  # pragma: no cover
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as exc:  # pragma: no cover
    raise EverythingNotInstalledError("Требуется установить пакет python-pptx") from exc

logger = logging.getLogger(__name__)

FILE_KIND_EXT = {
    "txt": ".txt",
    "docx": ".docx",
    "xlsx": ".xlsx",
    "pptx": ".pptx",
}

FILE_TYPE_EXT = {
    "текст": ".txt",
    "текстовый": ".txt",
    "txt": ".txt",
    "text": ".txt",
    "word": ".docx",
    "ворд": ".docx",
    "документ": ".docx",
    "docx": ".docx",
    "excel": ".xlsx",
    "таблица": ".xlsx",
    "xls": ".xlsx",
    "xlsx": ".xlsx",
    "ppt": ".pptx",
    "pptx": ".pptx",
    "презентация": ".pptx",
}

KIND_BY_EXTENSION = {value: key for key, value in FILE_KIND_EXT.items()}
FILE_KIND_ALIASES = {
    alias: KIND_BY_EXTENSION.get(ext)
    for alias, ext in FILE_TYPE_EXT.items()
    if ext in KIND_BY_EXTENSION
}

DOCX_CONTENT_TYPES_XML = """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
    <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
    <Default Extension="xml" ContentType="application/xml"/>
    <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
    <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
    <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
</Types>
"""

DOCX_RELS_XML = """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>
"""

DOCX_DOCUMENT_TEMPLATE = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas" xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" xmlns:o="urn:schemas-microsoft-com:office:office" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math" xmlns:v="urn:schemas-microsoft-com:vml" xmlns:wp14="http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing" xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing" xmlns:w10="urn:schemas-microsoft-com:office:word" xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:w14="http://schemas.microsoft.com/office/word/2010/wordml" xmlns:wpg="http://schemas.microsoft.com/office/word/2010/wordprocessingGroup" xmlns:wpi="http://schemas.microsoft.com/office/word/2010/wordprocessingInk" xmlns:wne="http://schemas.microsoft.com/office/word/2006/wordml" xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape" mc:Ignorable="w14 wp14">
  <w:body>
{paragraphs}
    <w:sectPr>
      <w:pgSz w:w="12240" w:h="15840"/>
      <w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440" w:header="708" w:footer="708" w:gutter="0"/>
      <w:cols w:space="708"/>
      <w:docGrid w:linePitch="360"/>
    </w:sectPr>
  </w:body>
</w:document>
"""

DOCX_DOCUMENT_RELS_XML = """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>
"""

DOCX_CORE_TEMPLATE = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
    <dc:title>Document</dc:title>
    <dc:creator>LocalWinAgent</dc:creator>
    <cp:lastModifiedBy>LocalWinAgent</cp:lastModifiedBy>
    <dcterms:created xsi:type="dcterms:W3CDTF">{timestamp}</dcterms:created>
    <dcterms:modified xsi:type="dcterms:W3CDTF">{timestamp}</dcterms:modified>
</cp:coreProperties>
"""

DOCX_APP_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
    <Application>LocalWinAgent</Application>
</Properties>
"""


def _expand(path: str | os.PathLike[str]) -> Path:
    expanded = os.path.expandvars(str(path))
    expanded = os.path.expanduser(expanded)
    return Path(expanded)


def normalize_path(
    path: str | os.PathLike[str], *, base: str | os.PathLike[str] | None = None
) -> Path:
    """Нормализовать путь до абсолютного представления."""

    candidate = _expand(path)
    if candidate.is_absolute():
        return candidate.resolve(strict=False)
    base_path = _expand(base) if base is not None else Path.cwd()
    if not base_path.is_absolute():
        base_path = (Path.cwd() / base_path).resolve(strict=False)
    return (base_path / candidate).resolve(strict=False)


def get_desktop_path() -> Path:
    desktop = config._KNOWN.get("DESKTOP")  # pylint: disable=protected-access
    if desktop:
        return Path(desktop).resolve(strict=False)
    return (Path.home() / "Desktop").resolve(strict=False)


def _get_file_attributes(path: Path) -> int:
    if platform.system() != "Windows":
        return 0
    try:
        import ctypes  # type: ignore

        attrs = ctypes.windll.kernel32.GetFileAttributesW(str(path))  # type: ignore[attr-defined]
    except Exception:  # pragma: no cover - системные ошибки
        return 0
    return attrs if attrs != -1 else 0


def _is_hidden_or_system(path: Path) -> bool:
    name = path.name
    if not name:
        return False
    lowered = name.lower()
    if lowered == "desktop.ini" or lowered.startswith("~$"):
        return True
    if name.startswith("."):
        return True
    attributes = _get_file_attributes(path)
    return bool(attributes & 0x2) or bool(attributes & 0x4)


def is_path_hidden(path: Path) -> bool:
    return _is_hidden_or_system(path)


def _resolve_shortcut(path: Path) -> Path:
    if platform.system() != "Windows" or path.suffix.lower() != ".lnk":
        return path
    try:  # pragma: no cover - зависит от win32com
        import win32com.client  # type: ignore

        shell = win32com.client.Dispatch("WScript.Shell")  # type: ignore[attr-defined]
        shortcut = shell.CreateShortCut(str(path))
        target = Path(shortcut.Targetpath)
        return target.resolve(strict=False)
    except Exception:  # pragma: no cover
        logger.warning("Не удалось разрешить ярлык %s", path)
        return path


def open_path(path: str | os.PathLike[str]) -> dict:
    manager = FileManager([])
    return manager.open_path(str(path))


class FileManager:
    """Класс для безопасной работы с файлами в заданных директориях."""

    def __init__(self, whitelist: Iterable[str]):
        self.whitelist = list(whitelist)
        self._allowed_paths = [normalize_path(item) for item in self.whitelist]
        self._default_root = (
            self._allowed_paths[0] if self._allowed_paths else get_desktop_path()
        )
        logger.debug("Белый список директорий: %s", self._allowed_paths)

    @property
    def default_root(self) -> Path:
        return self._default_root

    def normalize(self, path: str | os.PathLike[str]) -> Path:
        return normalize_path(path, base=self._default_root)

    def _normalize_kind(self, kind: Optional[str]) -> Optional[str]:
        if not kind:
            return None
        key = kind.strip().lower()
        if key in FILE_KIND_EXT:
            return key
        if key in FILE_KIND_ALIASES and FILE_KIND_ALIASES[key]:
            return FILE_KIND_ALIASES[key]
        return None

    def _resolve_path(self, raw: str, *, kind: Optional[str] = None) -> Path:
        candidate = _expand(raw)
        normalized_kind = self._normalize_kind(kind)
        if normalized_kind:
            ext = FILE_KIND_EXT[normalized_kind]
            if not candidate.suffix:
                candidate = candidate.with_name(candidate.name + ext)
        if not candidate.is_absolute():
            candidate = (self._default_root / candidate).resolve(strict=False)
        else:
            candidate = candidate.resolve(strict=False)
        return candidate

    def _is_allowed(self, path: Path) -> bool:
        for allowed in self._allowed_paths:
            try:
                path.relative_to(allowed)
                return True
            except ValueError:
                continue
        return False

    def _make_confirmation(self, path: Path, op_name: str) -> dict:
        logger.warning("Требуется подтверждение для %s: %s", op_name, path)
        exists, size = self._inspect(path)
        return self._finalize(
            path,
            ok=False,
            exists=exists,
            size=size,
            error=f"Требуется подтверждение для {op_name}",
            requires_confirmation=True,
            verified=False,
        )

    def _ensure_allowed(
        self, path: Path, confirmed: bool, op_name: str
    ) -> Optional[dict]:
        if self._is_allowed(path) or confirmed:
            return None
        return self._make_confirmation(path, op_name)

    def _inspect(self, path: Path) -> tuple[bool, int]:
        try:
            exists = path.exists()
            size = path.stat().st_size if exists and path.is_file() else 0
        except OSError:
            exists = False
            size = 0
        return exists, size

    def _finalize(
        self,
        path: Path,
        *,
        ok: bool,
        exists: Optional[bool] = None,
        size: Optional[int] = None,
        error: Optional[str] = None,
        requires_confirmation: bool = False,
        status: Optional[str] = None,
        verified: Optional[bool] = None,
    ) -> dict:
        result = {
            "ok": ok,
            "path": str(path),
            "requires_confirmation": requires_confirmation,
        }
        if exists is not None:
            result["exists"] = exists
        if size is not None:
            result["size"] = size
        if status is not None:
            result["status"] = status
        if verified is None:
            verified = ok and not requires_confirmation
        result["verified"] = verified
        if error:
            result["error"] = error
        return result

    def _handle_exception(self, path: Path, exc: Exception, op_name: str) -> dict:
        if isinstance(exc, NotAllowedPathError):
            error_message = str(exc)
        else:
            error_message = str(exc) or f"Ошибка {op_name}"
        logger.exception("Ошибка при %s %s: %s", op_name, path, exc)
        exists, size = self._inspect(path)
        return self._finalize(
            path,
            ok=False,
            exists=exists,
            size=size,
            error=error_message,
            verified=False,
        )

    def _detect_kind_from_path(self, path: Path, fallback: Optional[str] = None) -> str:
        ext = path.suffix.lower()
        if ext in KIND_BY_EXTENSION:
            return KIND_BY_EXTENSION[ext]
        normalized = self._normalize_kind(fallback)
        if normalized:
            return normalized
        return "txt"

    def _sync_write(self, path: Path, content: str, mode: str, encoding: str) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        with path.open(mode, encoding=encoding) as handler:
            handler.write(content)
            handler.flush()
            os.fsync(handler.fileno())

    def _prepare_docx_lines(self, content: Optional[str]) -> list[str]:
        if content is None:
            return []
        lines = content.splitlines()
        if not lines and content:
            lines = [content]
        return lines

    def _populate_new_document(self, document: "Document", lines: list[str]) -> None:  # type: ignore[name-defined]
        target_lines = lines or [""]
        if document.paragraphs:
            document.paragraphs[0].text = target_lines[0]
            remaining = target_lines[1:]
        else:  # pragma: no cover - защитный блок для нестандартных версий python-docx
            remaining = target_lines
            if target_lines:
                document.add_paragraph(target_lines[0])
                remaining = target_lines[1:]
        for line in remaining:
            document.add_paragraph(line)

    def _write_docx_fallback(self, target: Path, lines: list[str]) -> None:
        target.parent.mkdir(parents=True, exist_ok=True)
        safe_lines = lines or [""]
        paragraph_xml_parts = [
            "    <w:p>\n      <w:r>\n        <w:t xml:space=\"preserve\">{text}</w:t>\n      </w:r>\n    </w:p>".format(
                text=escape(line)
            )
            for line in safe_lines
        ]
        paragraphs_xml = "\n".join(paragraph_xml_parts)
        document_xml = DOCX_DOCUMENT_TEMPLATE.format(paragraphs=paragraphs_xml)
        timestamp = datetime.utcnow().replace(microsecond=0).isoformat() + "Z"
        core_xml = DOCX_CORE_TEMPLATE.format(timestamp=timestamp)
        with ZipFile(str(target), "w", ZIP_DEFLATED) as archive:
            archive.writestr("[Content_Types].xml", DOCX_CONTENT_TYPES_XML)
            archive.writestr("_rels/.rels", DOCX_RELS_XML)
            archive.writestr("word/document.xml", document_xml)
            archive.writestr("word/_rels/document.xml.rels", DOCX_DOCUMENT_RELS_XML)
            archive.writestr("docProps/core.xml", core_xml)
            archive.writestr("docProps/app.xml", DOCX_APP_XML)

    def _write_docx_document(self, target: Path, content: Optional[str]) -> None:
        lines = self._prepare_docx_lines(content)
        if Document is None:
            self._write_docx_fallback(target, lines)
            return
        target.parent.mkdir(parents=True, exist_ok=True)
        document = Document()
        self._populate_new_document(document, lines)
        document.save(str(target))

    def _append_docx_document(self, target: Path, lines: list[str], *, document, is_new: bool) -> None:
        if Document is None:
            self._write_docx_fallback(target, lines)
            return
        if is_new:
            self._populate_new_document(document, lines)
        else:
            for line in lines:
                document.add_paragraph(line)
        document.save(str(target))

    def create_file(
        self,
        path: str,
        content: Optional[str] = None,
        *,
        kind: Optional[str] = None,
        confirmed: bool = False,
        encoding: str = "utf-8",
    ) -> dict:
        target = self._resolve_path(path, kind=kind)
        denied = self._ensure_allowed(target, confirmed, "создания файла")
        if denied:
            return denied
        logger.info("Создание файла: %s", target)
        existed_before = target.exists()
        file_kind = self._detect_kind_from_path(target, kind)
        try:
            if file_kind == "docx":
                self._write_docx_document(target, content)
            elif file_kind == "xlsx":
                target.parent.mkdir(parents=True, exist_ok=True)
                workbook = Workbook()
                sheet = workbook.active
                if content:
                    sheet["A1"] = str(content)
                workbook.save(str(target))
            elif file_kind == "pptx":
                target.parent.mkdir(parents=True, exist_ok=True)
                presentation = Presentation()
                layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
                slide = presentation.slides.add_slide(layout)
                if content:
                    applied = False
                    title = getattr(slide.shapes, "title", None)
                    if title is not None and getattr(title, "has_text_frame", False):
                        title.text = content
                        applied = True
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            shape.text_frame.text = content
                            applied = True
                            break
                    if not applied:
                        textbox = slide.shapes.add_textbox(
                            Inches(1), Inches(1.5), Inches(8), Inches(2)
                        )
                        textbox.text_frame.text = content
                presentation.save(str(target))
            else:
                text = content if content is not None else ""
                self._sync_write(target, text, mode="w", encoding=encoding)
            exists, size = self._inspect(target)
            logger.info("Создание завершено: %s exists=%s size=%s", target, exists, size)
            status = "updated" if existed_before else "created"
            return self._finalize(target, ok=exists, exists=exists, size=size, status=status)
        except Exception as exc:
            return self._handle_exception(target, exc, "создании файла")

    def write_text(
        self,
        path: str,
        content: str,
        *,
        confirmed: bool = False,
        encoding: str = "utf-8",
    ) -> dict:
        target = self._resolve_path(path)
        denied = self._ensure_allowed(target, confirmed, "записи файла")
        if denied:
            return denied
        logger.info("Запись текста в файл: %s", target)
        existed_before = target.exists()
        try:
            self._sync_write(target, content, mode="w", encoding=encoding)
            exists, size = self._inspect(target)
            logger.info("Запись завершена: %s exists=%s size=%s", target, exists, size)
            status = "overwritten" if existed_before else "created"
            return self._finalize(target, ok=exists, exists=exists, size=size, status=status)
        except Exception as exc:
            return self._handle_exception(target, exc, "записи файла")

    def append_text(
        self,
        path: str,
        content: str,
        *,
        confirmed: bool = False,
        encoding: str = "utf-8",
    ) -> dict:
        target = self._resolve_path(path)
        denied = self._ensure_allowed(target, confirmed, "добавления текста")
        if denied:
            return denied
        logger.info("Добавление текста в файл: %s", target)
        existed_before = target.exists()
        try:
            mode = "a" if target.exists() else "w"
            self._sync_write(target, content, mode=mode, encoding=encoding)
            exists, size = self._inspect(target)
            logger.info("Добавление завершено: %s exists=%s size=%s", target, exists, size)
            status = "appended" if existed_before else "created"
            return self._finalize(target, ok=exists, exists=exists, size=size, status=status)
        except Exception as exc:
            return self._handle_exception(target, exc, "добавлении текста")

    def read_text(self, path: str, encoding: str = "utf-8") -> dict:
        target = self._resolve_path(path)
        logger.info("Чтение файла: %s", target)
        if target.suffix.lower() != ".txt":
            error = "Чтение доступно только для текстовых файлов"
            logger.error("%s: %s", error, target)
            return self._finalize(target, ok=False, exists=target.exists(), error=error)
        try:
            content = target.read_text(encoding=encoding)
            exists, size = self._inspect(target)
            logger.info("Чтение завершено: %s exists=%s size=%s", target, exists, size)
            return {
                "ok": True,
                "path": str(target),
                "content": content,
                "exists": exists,
                "size": size,
                "requires_confirmation": False,
            }
        except Exception as exc:
            return self._handle_exception(target, exc, "чтении файла")

    def edit_word(
        self,
        path: str,
        content: str,
        *,
        confirmed: bool = False,
    ) -> dict:
        target = self._resolve_path(path, kind="docx")
        denied = self._ensure_allowed(target, confirmed, "редактирования документа")
        if denied:
            return denied
        logger.info("Редактирование Word-документа: %s", target)
        existed_before = target.exists()
        lines = self._prepare_docx_lines(content)
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            if Document is None:
                if existed_before:
                    logger.warning(
                        "Библиотека python-docx недоступна, пересоздаём документ %s в упрощённом формате.",
                        target,
                    )
                self._write_docx_fallback(target, lines)
            else:
                document = None
                is_new_document = not existed_before
                if existed_before:
                    try:
                        document = Document(str(target))
                        is_new_document = False
                    except (PackageNotFoundError, OSError, ValueError, KeyError) as exc:
                        logger.warning(
                            "Файл %s повреждён или недоступен, пересоздаём документ: %s",
                            target,
                            exc,
                        )
                        document = Document()
                        is_new_document = True
                if document is None:
                    document = Document()
                    is_new_document = True
                self._append_docx_document(target, lines, document=document, is_new=is_new_document)
            exists, size = self._inspect(target)
            logger.info("Word-документ обновлён: %s exists=%s size=%s", target, exists, size)
            status = "updated" if existed_before else "created"
            return self._finalize(target, ok=exists, exists=exists, size=size, status=status)
        except Exception as exc:
            return self._handle_exception(target, exc, "редактировании документа")

    def edit_excel(
        self,
        path: str,
        cell: str,
        value: str,
        *,
        confirmed: bool = False,
    ) -> dict:
        target = self._resolve_path(path, kind="xlsx")
        denied = self._ensure_allowed(target, confirmed, "редактирования таблицы")
        if denied:
            return denied
        logger.info("Редактирование Excel-файла: %s", target)
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            if target.exists():
                workbook = load_workbook(str(target))
            else:
                workbook = Workbook()
            sheet = workbook.active
            sheet[str(cell).upper()] = value
            workbook.save(str(target))
            exists, size = self._inspect(target)
            logger.info("Excel-файл обновлён: %s exists=%s size=%s", target, exists, size)
            return self._finalize(target, ok=exists, exists=exists, size=size)
        except Exception as exc:
            return self._handle_exception(target, exc, "редактировании таблицы")

    def edit_pptx(
        self,
        path: str,
        content: str,
        *,
        confirmed: bool = False,
    ) -> dict:
        target = self._resolve_path(path, kind="pptx")
        denied = self._ensure_allowed(target, confirmed, "редактирования презентации")
        if denied:
            return denied
        logger.info("Редактирование презентации: %s", target)
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            presentation = Presentation(str(target)) if target.exists() else Presentation()
            layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(layout)
            if content:
                applied = False
                title = getattr(slide.shapes, "title", None)
                if title is not None and getattr(title, "has_text_frame", False):
                    title.text = content
                    applied = True
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        shape.text_frame.text = content
                        applied = True
                        break
                if not applied:
                    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
                    textbox.text_frame.text = content
            presentation.save(str(target))
            exists, size = self._inspect(target)
            logger.info("Презентация обновлена: %s exists=%s size=%s", target, exists, size)
            return self._finalize(target, ok=exists, exists=exists, size=size)
        except Exception as exc:
            return self._handle_exception(target, exc, "редактировании презентации")

    def open_path(self, path: str) -> dict:
        target = self._resolve_path(path)
        logger.info("Открытие пути: %s", target)
        if not target.exists():
            logger.error("Путь не существует: %s", target)
            return self._finalize(target, ok=False, exists=False, error="Путь не существует")
        actual = _resolve_shortcut(target)
        resolved = actual.resolve(strict=False)
        try:
            if os.name == "nt":
                os.startfile(str(actual))  # type: ignore[attr-defined]
            else:
                if actual.is_dir():
                    subprocess.Popen(["xdg-open", str(actual)])  # noqa: S603,S607 pragma: no cover - для *nix
                else:
                    open_with_shell(str(actual))
        except Exception as exc:  # pragma: no cover - системные ошибки
            logger.exception("Ошибка открытия %s: %s", resolved, exc)
            exists, size = self._inspect(resolved)
            return self._finalize(resolved, ok=False, exists=exists, size=size, error=str(exc))
        logger.info("Путь открыт: %s", resolved)
        exists, size = self._inspect(resolved)
        return self._finalize(resolved, ok=True, exists=exists, size=size)

    def move_path(self, src: str, dst: str, *, confirmed: bool = False) -> dict:
        source = self._resolve_path(src)
        destination = self._resolve_path(dst)
        denied_src = self._ensure_allowed(source, confirmed, "перемещения файла")
        if denied_src:
            return denied_src
        denied_dst = self._ensure_allowed(destination, confirmed, "перемещения файла")
        if denied_dst:
            return denied_dst
        logger.info("Перемещение: %s -> %s", source, destination)
        try:
            destination.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(source), str(destination))
            exists, size = self._inspect(destination)
            logger.info("Перемещение завершено: %s exists=%s size=%s", destination, exists, size)
            return self._finalize(destination, ok=exists, exists=exists, size=size)
        except Exception as exc:
            return self._handle_exception(destination, exc, "перемещении файла")

    def copy_path(self, src: str, dst: str, *, confirmed: bool = False) -> dict:
        source = self._resolve_path(src)
        destination = self._resolve_path(dst)
        denied_src = self._ensure_allowed(source, confirmed, "копирования файла")
        if denied_src:
            return denied_src
        denied_dst = self._ensure_allowed(destination, confirmed, "копирования файла")
        if denied_dst:
            return denied_dst
        logger.info("Копирование: %s -> %s", source, destination)
        try:
            if source.is_dir():
                shutil.copytree(source, destination, dirs_exist_ok=True)
            else:
                destination.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(source, destination)
            exists, size = self._inspect(destination)
            logger.info("Копирование завершено: %s exists=%s size=%s", destination, exists, size)
            return self._finalize(destination, ok=exists, exists=exists, size=size)
        except Exception as exc:
            return self._handle_exception(destination, exc, "копировании файла")

    def delete_path(self, path: str, *, confirmed: bool = False) -> dict:
        target = self._resolve_path(path)
        denied = self._ensure_allowed(target, confirmed, "удаления файла")
        if denied:
            return denied
        logger.info("Удаление: %s", target)
        try:
            if target.is_dir():
                shutil.rmtree(target, ignore_errors=False)
            elif target.exists():
                target.unlink()
            exists, _ = self._inspect(target)
            logger.info("Удаление завершено: %s exists=%s", target, exists)
            return self._finalize(target, ok=not exists, exists=exists)
        except Exception as exc:
            return self._handle_exception(target, exc, "удалении файла")

    def list_directory(self, path: Optional[str] = None, *, confirmed: bool = False) -> dict:
        directory = self._default_root if path is None else self._resolve_path(path)
        denied = self._ensure_allowed(directory, confirmed, "просмотра каталога")
        if denied:
            return denied
        logger.info("Просмотр каталога: %s", directory)
        try:
            if not directory.exists() or not directory.is_dir():
                raise FileOperationError(f"Каталог {directory} не найден")
            items = [
                item.name
                for item in sorted(directory.iterdir(), key=lambda candidate: candidate.name.lower())
                if not is_path_hidden(item)
            ]
            logger.info("Каталог %s содержит %d элементов", directory, len(items))
            return {
                "ok": True,
                "path": str(directory),
                "items": items,
                "requires_confirmation": False,
            }
        except Exception as exc:
            return self._handle_exception(directory, exc, "просмотре каталога")
